from models.get_transcript import get_vectorstore
from commands.command import Command
from langchain_openai.chat_models import ChatOpenAI
from schemas.pydantic_models import QueryResponse, QueryInput, SlideResponse, SlideSchema
from langchain_core.output_parsers import PydanticOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from dotenv import load_dotenv
import os
from typing import List

load_dotenv()

class GeneratePPT(Command):

        
    def __init__(self, input:QueryInput):
        self.session_id = input.session_id
        self.model_name = input.model
        self.query = input.query

    def save_presentation(self, presentation):
        try:
            os.makedirs('presentations', exist_ok=True)
            
            filename = os.path.join('presentations', f'{self.session_id}.pptx')
            presentation.save(filename)
            
            print(f"Presentation saved to: {filename}")
            return filename
            
        except Exception as e:
            print(f"Error saving presentation: {e}")
            return None
        
    def execute(self):
        vectorstore = get_vectorstore(self.session_id)
        qa_response = run_langchain_qa_chain(vectorstore, self.query, self.model_name, self.session_id)

        prs = Presentation()
        title = prs.slides.add_slide(prs.slide_layouts[0])
        title.shapes.title.text = f"Answering: {self.query}"
        title.placeholders[1].text = "Created using sunpters' website"

        slides: SlideResponse = qa_response.response
        i = 1
        for slide_schema in slides.content:
            new_slide = prs.slides.add_slide(prs.slide_layouts[6]) #blank slide
            slide_schema: SlideSchema
            slide_template(new_slide, slide_schema.title, slide_schema.slide_content, slide_schema.bullet_points)
            print(f"Generating slide {i}")
            i+=1
            
        self.save_presentation(prs)

          
'''
Format the slide according to a predetermined layout for the MVP
'''
def slide_template(slide, title: str, content: str, bullets: List[str], is_image=False):
#According to the pydantic model for SlideSchema 
#Generate image later when I buy more tokens :)

    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1)) #Styling help from claude
    header_frame = header_box.text_frame
    header_frame.text = title 

    '''
    Styling for the header
    '''
    header = header_frame.paragraphs[0]
    header.alignment = PP_ALIGN.LEFT
    header.font.name = 'Arial'
    header.font.size = Pt(28)
    header.font.bold = True

    '''
    Add the bullet points
    '''
    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(5))
    bullet_text = bullet_box.text_frame
    bullet_text.word_wrap = True

    for i, bullet_point in enumerate(bullets):
        if i == 0:
            # First paragraph is already created
            p = bullet_text.paragraphs[0]
        else:
            # Add new paragraphs for additional bullet points
            p = bullet_text.add_paragraph()
        
        p.text = bullet_point
        p.level = 0  # Main bullet level
        
        # Format bullet points
        font = p.font
        font.name = 'Arial'
        font.size = Pt(16)
        font.color.rgb = RGBColor(51, 51, 51)  # Dark gray
    
    '''Content Text'''
    context_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.3), Inches(5))
    context_frame = context_box.text_frame
    context_frame.word_wrap = True

    context_frame.text = content
    context_font = context_frame.paragraphs[0].font
    context_font.name = 'Arial'
    context_font.size = Pt(14)
    
    return slide 

def run_langchain_qa_chain(vectorstore, prompt: str, model_name: str, session_id:str) -> QueryResponse:
    retriever = vectorstore.as_retriever(search_kwargs={"k": 5})
    llm = ChatOpenAI(model_name=model_name)

    parser = PydanticOutputParser(pydantic_object=SlideResponse)
    format_instructions = parser.get_format_instructions()

    PROMPT_TEMPLATE = ChatPromptTemplate.from_template(
    """Generate a comprehensive presentation based on the following context and question.
        Create exactly 10 slides that thoroughly cover the topic.

        Context: {context}

        Question: {question}

        {format_instructions}

        Generate slides that cover the topic comprehensively. Each slide should have:
        - A clear, descriptive title
        - 3-5 bullet points with key information
        - Detailed slide content that expands on the bullet points
        - A boolean indicating if an image would enhance understanding

        Response:"""
    )

    qa_chain = (
        {
            "context": retriever,
            "question": RunnablePassthrough(),
            "format_instructions": lambda _: format_instructions
        }
        | PROMPT_TEMPLATE
        | llm
        | parser
    )

    try: 
        response = qa_chain.invoke(prompt) 
        
        return QueryResponse (
            response = response,
            session_id=session_id,
            model=model_name
        )
        
    except Exception as e:
        print(f"Raised exception {e}")
        return None

